"""
批量导入飞书文档到三层数据架构 — 使用 export + python-pptx 提取 slides 文本
"""
import json
import os
import sys
import subprocess
import hashlib
import time
from pathlib import Path

sys.path.insert(0, "/workspace/services/orchestrator")

RAW_DIR = Path("/data/raw/feishu")
EXPORT_DIR = Path("/tmp/feishu_exports")
EXPORT_DIR.mkdir(parents=True, exist_ok=True)

# 完整文档清单
DOCUMENTS = [
    # (token, type, name, source_folder)
    ("BHmksvrNVlZjYod18igc0ZeSnUc", "slides", "英氏_KOS_全链路新客转化方案_0609", "根目录"),
    ("OgDBsf53Bl7xqmdQmPwcuJZXnQg", "slides", "金领冠小红书KOS营销方案_0616", "根目录"),
    ("AD2TsNeU4lElVSdp6KbcWna4nHc", "slides", "飞鹤KOS搜索占位与评论区营销方案0605", "根目录"),
    ("CLVNsgBfvlhXhtdMmqRcpUIOnHe", "slides", "派特生物KOS推广方案_0923_v1", "根目录"),
    ("Mw99dLUVZoXYTBxl7DFcGxRCntH", "docx", "繁星计划AI_Service解决方案", "根目录"),
    ("XJfQsq6MZl976GdsK68cDnTCnZb", "slides", "繁星计划AI_Service解决方案_slides", "根目录"),
    ("EDXOsmIC5lDF0adcJWtcYb5Anob", "slides", "火山引擎_Agent_DataLake解决方案", "根目录"),
    ("SVU4sFra6l9t3tdzxXYcU4HCndc", "slides", "母婴行业小红书KOS矩阵营销通案", "行业通案"),
    ("DfLXsxHoylzvE2dgaGScmCPDnBb", "slides", "大健康行业小红书KOS矩阵营销通案", "行业通案"),
    ("Qz3UsoeY8lNVludxgmGc1jyLnFe", "slides", "0610繁星计划提质提效解决方案", "飞鹤"),
    ("W0SGsB97nluTuHdNKrmcXaGMn0c", "slides", "飞鹤繁星计划升级方案", "飞鹤"),
    ("Mah3sAY1WlqnzzdxRijcgoVAn8c", "slides", "飞鹤启萃_KOS全链路新客转化方案", "飞鹤"),
    ("Z9rKswa0plcOw6dX18pc7YlqnOc", "slides", "终版_飞鹤繁星计划koc孵化运营方案v8", "飞鹤"),
]


def fetch_docx(token, name):
    """读取 docx 文档"""
    url = f"https://ycnm3444stv0.feishu.cn/docx/{token}"
    result = subprocess.run(
        ["lark-cli", "docs", "+fetch", "--api-version", "v2", "--doc", url,
         "--doc-format", "markdown", "--as", "user"],
        capture_output=True, text=True, timeout=90
    )
    if result.returncode != 0 or not result.stdout.strip():
        return None
    return result.stdout


def export_and_extract_slides(token, name):
    """导出 slides 为 pptx 并提取文本"""
    safe_name = name.replace(" ", "_").replace("/", "_")[:60]
    pptx_path = EXPORT_DIR / f"{safe_name}.pptx"
    
    # 如果已导出，直接使用
    if pptx_path.exists():
        return extract_pptx_text(str(pptx_path))
    
    # 导出 slides 为 pptx
    result = subprocess.run(
        ["lark-cli", "drive", "+export",
         "--token", token,
         "--doc-type", "slides",
         "--file-extension", "pptx",
         "--file-name", str(pptx_path),
         "--overwrite",
         "--as", "user"],
        capture_output=True, text=True, timeout=120
    )
    
    # 等待导出完成
    time.sleep(3)
    
    if pptx_path.exists():
        return extract_pptx_text(str(pptx_path))
    
    # 如果 timeout，尝试查 task_result
    if "timed_out" in result.stdout or "timed_out" in result.stderr:
        try:
            data = json.loads(result.stdout)
            ticket = data.get("ticket", "")
            if ticket:
                time.sleep(10)
                subprocess.run(
                    ["lark-cli", "drive", "+task_result", "--scenario", "export",
                     "--ticket", ticket, "--file-token", token, "--as", "user"],
                    capture_output=True, text=True, timeout=60
                )
                time.sleep(3)
                if pptx_path.exists():
                    return extract_pptx_text(str(pptx_path))
        except:
            pass
    
    return None


def extract_pptx_text(pptx_path):
    """从 PPTX 提取文本"""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        texts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t and len(t) > 1:
                            slide_texts.append(t)
            if slide_texts:
                texts.append(f"## 第{i+1}页\n" + "\n".join(slide_texts))
        return "\n\n".join(texts)
    except Exception as e:
        return None


def classify_document(name, content):
    """根据文档名称分类到知识库"""
    classifications = []
    n = name
    
    if any(w in n for w in ["行业", "通案", "矩阵"]):
        classifications.append("industry_strategy")
    if any(w in n for w in ["竞品", "对比", "对标"]):
        classifications.append("competitor_analysis")
    if any(w in n for w in ["繁星", "转化", "增长", "KOS", "全链路", "孵化"]):
        classifications.append("growth_model")
    if any(w in n for w in ["AI Service", "Agent", "数据", "内容罗盘", "DataLake","解决方案"]):
        classifications.append("product_solution")
    if any(w in n for w in ["案例", "复盘", "onepage", "代管代发"]):
        classifications.append("case_labels")
    if any(w in n for w in ["飞鹤", "英氏", "金领冠", "派特"]):
        classifications.append("brand_knowledge")
    if any(w in n for w in ["方案", "推广", "营销", "计划"]):
        classifications.append("proposal_review")
    
    return list(set(classifications)) or ["industry_strategy"]


def extract_knowledge(name, content, classifications):
    """从文档提取知识点"""
    knowledge = {"collections": {}, "relations": []}
    
    if not content:
        return knowledge
    
    lines = content.split("\n")
    headings = []
    for l in lines:
        l = l.strip()
        if l.startswith("## ") or l.startswith("# "):
            h = l.lstrip("# ").strip()
            if len(h) > 3:
                headings.append(h)
    
    for coll in classifications:
        items = []
        for h in headings[:15]:
            items.append({"title": h, "summary": f"来源: {name}"})
        if items:
            knowledge["collections"][coll] = items
    
    # 建立关联
    doc_hash = hashlib.md5(name.encode()).hexdigest()[:12]
    for i, c1 in enumerate(classifications):
        for c2 in classifications[i+1:]:
            knowledge["relations"].append({
                "source_type": c1, "source_id": doc_hash,
                "target_type": c2, "target_id": f"{doc_hash}_{c2}",
                "relation_type": "belongs_to", "weight": 0.8,
            })
    
    return knowledge


def main():
    RAW_DIR.mkdir(parents=True, exist_ok=True)
    
    stats = {"docx": 0, "slides": 0, "failed": 0}
    all_knowledge = {"collections": {}, "relations": []}
    
    for token, doc_type, name, folder in DOCUMENTS:
        print(f"\n{'='*50}")
        print(f"[{doc_type}] {name}")
        print(f"  文件夹: {folder}")
        
        safe_name = name.replace(" ", "_").replace("/", "_")[:80]
        
        # 获取内容
        content = None
        if doc_type == "docx":
            content = fetch_docx(token, name)
            if content:
                stats["docx"] += 1
        elif doc_type == "slides":
            content = export_and_extract_slides(token, name)
            if content:
                stats["slides"] += 1
        
        if not content:
            print(f"  [SKIP] 无法获取内容")
            stats["failed"] += 1
            continue
        
        print(f"  内容长度: {len(content)} 字符")
        
        # L1: 保存到文件系统
        doc_path = RAW_DIR / folder / f"{safe_name}.md"
        doc_path.parent.mkdir(parents=True, exist_ok=True)
        doc_path.write_text(content, encoding="utf-8")
        print(f"  L1 -> {doc_path}")
        
        # 分类
        classifications = classify_document(name, content)
        print(f"  分类: {classifications}")
        
        # 提取知识
        knowledge = extract_knowledge(name, content, classifications)
        for coll, items in knowledge["collections"].items():
            if coll not in all_knowledge["collections"]:
                all_knowledge["collections"][coll] = []
            all_knowledge["collections"][coll].extend(items)
            print(f"  L2 -> {coll}: +{len(items)} 条")
        all_knowledge["relations"].extend(knowledge["relations"])
        
        # 保存知识 JSON
        knowledge_path = RAW_DIR / folder / f"{safe_name}_knowledge.json"
        knowledge_path.write_text(json.dumps(knowledge, ensure_ascii=False, indent=2))
    
    # 汇总知识
    summary_path = RAW_DIR / "knowledge_summary.json"
    summary_path.write_text(json.dumps(all_knowledge, ensure_ascii=False, indent=2))
    
    print(f"\n{'='*50}")
    print(f"导入完成！")
    print(f"  docx: {stats['docx']} | slides: {stats['slides']} | 失败: {stats['failed']}")
    total_kp = sum(len(v) for v in all_knowledge["collections"].values())
    print(f"  知识点总数: {total_kp}")
    print(f"  关联数: {len(all_knowledge['relations'])}")
    print(f"  存储路径: {RAW_DIR}")
    print(f"  汇总文件: {summary_path}")


if __name__ == "__main__":
    main()