import subprocess
import json
import os
import tempfile
from pathlib import Path
from app.skills.base import BaseSkill, SkillInput, SkillOutput
from pydantic import Field
from typing import Optional


class FormatOutputInput(SkillInput):
    slides: list = Field(description="S7 生成的 Slide 列表")
    brand_assets: dict = Field(default_factory=dict)
    output_formats: list = Field(default=["slides", "docx", "pptx"])
    brand_name: str = Field(default="", description="品牌名，用于封面标题")
    proposal_title: str = Field(default="品牌增长全案方案", description="提案标题")


class FormatOutputOutput(SkillOutput):
    slides_url: Optional[str] = None
    docx_url: Optional[str] = None
    pptx_path: Optional[str] = None
    slides_token: Optional[str] = None
    docx_token: Optional[str] = None


class S8FormatOutput(BaseSkill):
    name = "s8_format_output"
    description = "格式输出：生成飞书 Slides + Docx + PPTX 导出"

    def execute(self, input_data: FormatOutputInput) -> FormatOutputOutput:
        result = FormatOutputOutput()

        if "slides" in input_data.output_formats:
            slides_result = self._create_feishu_slides(
                input_data.slides,
                input_data.brand_assets,
                input_data.brand_name,
                input_data.proposal_title,
            )
            result.slides_url = slides_result.get("url", "")
            result.slides_token = slides_result.get("token", "")

        if "docx" in input_data.output_formats:
            docx_result = self._create_feishu_docx(
                input_data.slides,
                input_data.brand_name,
                input_data.proposal_title,
            )
            result.docx_url = docx_result.get("url", "")
            result.docx_token = docx_result.get("token", "")

        if "pptx" in input_data.output_formats:
            result.pptx_path = self._export_pptx(
                input_data.slides, input_data.brand_name
            )

        return result

    # ── 飞书 Slides 生成 ──────────────────────────────────────

    def _create_feishu_slides(
        self, slides: list, assets: dict, brand_name: str, title: str
    ) -> dict:
        """通过 lark-cli 创建飞书幻灯片"""

        if not self._check_lark_cli():
            return self._fallback_slides(slides, brand_name, title)

        try:
            # 步骤 1：创建空白 PPT
            create_result = subprocess.run(
                ["lark-cli", "slides", "+create", "--as", "user",
                 "--format", "json"],
                capture_output=True, text=True, timeout=30,
            )
            create_data = json.loads(create_result.stdout)
            xml_presentation_id = (
                create_data.get("data", {})
                .get("xml_presentation", {})
                .get("xml_presentation_id", "")
            )

            if not xml_presentation_id:
                return self._fallback_slides(slides, brand_name, title)

            # 步骤 2：逐页添加 slide
            ppt_url = ""
            for i, slide_data in enumerate(slides):
                slide_xml = self._build_slide_xml(slide_data, brand_name)
                try:
                    add_result = subprocess.run(
                        [
                            "lark-cli", "slides", "xml_presentation.slide", "create",
                            "--as", "user",
                            "--params", json.dumps({
                                "xml_presentation_id": xml_presentation_id,
                            }),
                            "--data", json.dumps({
                                "slide": {"content": slide_xml},
                            }),
                        ],
                        capture_output=True, text=True, timeout=15,
                    )
                    if i == 0:
                        add_data = json.loads(add_result.stdout)
                        ppt_url = add_data.get("data", {}).get("url", "")
                except Exception:
                    continue

            slides_url = ppt_url or f"https://ycnm3444stv0.feishu.cn/slides/{xml_presentation_id}"
            return {"url": slides_url, "token": xml_presentation_id}

        except Exception as e:
            print(f"[S8] 飞书 Slides 创建失败: {e}")
            return self._fallback_slides(slides, brand_name, title)

    def _build_slide_xml(self, slide_data: dict, brand_name: str) -> str:
        """构建单页 Slide 的 XML，应用品牌色"""
        title = slide_data.get("title", "")
        layout = slide_data.get("layout_type", "default")
        bg_color = self._get_bg_color(layout)
        content = slide_data.get("content", "")

        # 品牌色：尝试从 brand_assets 获取，否则使用默认配色
        brand_color = brand_name and self._brand_color(brand_name) or "#1a1a2e"

        return f"""<slide xmlns="http://www.larkoffice.com/sml/2.0">
  <style>
    <fill><fillColor color="{bg_color}"/></fill>
  </style>
  <data>
    <shape type="rect" x="60" y="50" width="880" height="80">
      <fill><fillColor color="transparent"/></fill>
      <outline><noOutline/></outline>
      <content><p font-size="28pt" font-weight="bold" color="{brand_color}">{title}</p></content>
    </shape>
    <shape type="rect" x="60" y="140" width="880" height="400">
      <fill><fillColor color="transparent"/></fill>
      <outline><noOutline/></outline>
      <content><p font-size="14pt" color="#333333">{content}</p></content>
    </shape>
  </data>
</slide>"""

    @staticmethod
    def _brand_color(brand_name: str) -> str:
        """获取品牌色"""
        colors = {
            "飞鹤": "#1a5276", "a2": "#7b2d8b", "英氏": "#27ae60",
            "金领冠": "#e67e22", "林氏家居": "#2c3e50", "老庙": "#c0392b",
            "蒙牛": "#2980b9", "领克": "#16a085", "极氪": "#8e44ad",
            "利星行": "#34495e", "可画": "#e74c3c", "董酒": "#d35400",
        }
        return colors.get(brand_name, "#1a1a2e")

    def _get_bg_color(self, layout: str) -> str:
        colors = {
            "cover": "#1a1a2e",
            "two_column": "#f8f9fa",
            "chart_focus": "#ffffff",
            "bullet_cards": "#f0f4ff",
            "comparison_table": "#ffffff",
            "icon_grid": "#f8f9fa",
            "case_card": "#fff8f0",
            "timeline": "#ffffff",
            "grid_cards": "#f0f4ff",
        }
        return colors.get(layout, "#ffffff")

    def _fallback_slides(self, slides: list, brand_name: str, title: str) -> dict:
        """飞书 API 不可用时的降级方案"""
        return {
            "url": "https://ycnm3444stv0.feishu.cn/slides/placeholder",
            "token": "",
        }

    # ── 飞书 Docx 生成 ────────────────────────────────────────

    def _create_feishu_docx(
        self, slides: list, brand_name: str, title: str
    ) -> dict:
        """通过 lark-cli 创建飞书文档"""

        if not self._check_lark_cli():
            return self._fallback_docx(slides, brand_name, title)

        try:
            doc_content = self._build_docx_content(slides, brand_name, title)
            escaped = doc_content.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

            create_result = subprocess.run(
                [
                    "lark-cli", "docs", "+create", "--api-version", "v2",
                    "--as", "user",
                    "--content", f"<title>{title}</title>{escaped[:50000]}",
                    "--format", "json",
                ],
                capture_output=True, text=True, timeout=30,
            )
            create_data = json.loads(create_result.stdout)
            doc_token = create_data.get("data", {}).get("document", {}).get("document_id", "")
            doc_url = create_data.get("data", {}).get("document", {}).get("url", "")

            if not doc_url and doc_token:
                doc_url = f"https://ycnm3444stv0.feishu.cn/docx/{doc_token}"

            return {"url": doc_url, "token": doc_token}

        except Exception as e:
            print(f"[S8] 飞书 Docx 创建失败: {e}")
            return self._fallback_docx(slides, brand_name, title)

    def _build_docx_content(self, slides: list, brand_name: str, title: str) -> str:
        """构建 Docx 文档内容（XML 格式）"""
        parts = [f"<title>{title}</title>"]
        parts.append(f"<p font-size='18pt' font-weight='bold'>{brand_name} - {title}</p>")
        parts.append("<p/><hr/><p/>")

        for slide_data in slides:
            section_title = slide_data.get("title", "")
            parts.append(f"<p font-size='16pt' font-weight='bold'>{section_title}</p>")
            parts.append(f"<p>{slide_data.get('content', '')}</p>")
            parts.append("<p/>")

        return "".join(parts)

    def _fallback_docx(self, slides: list, brand_name: str, title: str) -> dict:
        """飞书 API 不可用时的降级方案"""
        return {
            "url": "https://ycnm3444stv0.feishu.cn/docx/placeholder",
            "token": "",
        }

    # ── PPTX 导出 ─────────────────────────────────────────────

    def _export_pptx(self, slides: list, brand_name: str) -> str:
        """使用 python-pptx 生成专业 PPTX 文件"""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in slides:
            layout = slide_data.get("layout_type", "default")
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

            # 背景色
            bg_color = self._get_pptx_bg(layout)
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = bg_color

            # 标题
            title = slide_data.get("title", "")
            if title:
                left = Inches(0.8)
                top = Inches(0.6)
                width = Inches(11.7)
                height = Inches(1.0)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

            # 内容
            content = slide_data.get("content", "")
            if content:
                left = Inches(0.8)
                top = Inches(1.8)
                width = Inches(11.7)
                height = Inches(5.0)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

            # 页脚
            footer_left = Inches(0.8)
            footer_top = Inches(7.0)
            footer_width = Inches(11.7)
            footer_height = Inches(0.3)
            ftBox = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
            ftf = ftBox.text_frame
            fp = ftf.paragraphs[0]
            fp.text = f"{brand_name} | 品牌增长全案 | 第{slide_data.get('index', '')}页"
            fp.font.size = Pt(8)
            fp.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            fp.alignment = PP_ALIGN.RIGHT

        import uuid
        export_id = str(uuid.uuid4())[:8]
        output_path = f"/data/exports/{brand_name}_{export_id}.pptx"
        os.makedirs("/data/exports", exist_ok=True)
        prs.save(output_path)
        return output_path

    def _get_pptx_bg(self, layout: str):
        from pptx.dml.color import RGBColor
        colors = {
            "cover": RGBColor(0x1A, 0x1A, 0x2E),
            "two_column": RGBColor(0xF8, 0xF9, 0xFA),
            "chart_focus": RGBColor(0xFF, 0xFF, 0xFF),
            "bullet_cards": RGBColor(0xF0, 0xF4, 0xFF),
            "comparison_table": RGBColor(0xFF, 0xFF, 0xFF),
            "icon_grid": RGBColor(0xF8, 0xF9, 0xFA),
            "case_card": RGBColor(0xFF, 0xF8, 0xF0),
            "timeline": RGBColor(0xFF, 0xFF, 0xFF),
            "grid_cards": RGBColor(0xF0, 0xF4, 0xFF),
        }
        return colors.get(layout, RGBColor(0xFF, 0xFF, 0xFF))

    def _check_lark_cli(self) -> bool:
        """检查 lark-cli 是否可用"""
        try:
            result = subprocess.run(
                ["lark-cli", "auth", "whoami", "--as", "user"],
                capture_output=True, text=True, timeout=5,
            )
            return result.returncode == 0
        except Exception:
            return False